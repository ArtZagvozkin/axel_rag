from __future__ import annotations

import io
from typing import Optional, List

from config import logger
from .base import FileTextExtractor, register_extractor

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


class PptxExtractor(FileTextExtractor):
    PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def supports(self, *, ext: str, mime_type: str | None) -> bool:
        mt = (mime_type or "").lower()
        return ext == "pptx" or mt == self.PPTX_MIME

    def extract_text(self, data: bytes, name: str | None, mime_type: str | None) -> Optional[str]:
        if Presentation is None:
            logger.warning("python-pptx is not installed; cannot extract PPTX text")
            return None

        prs = Presentation(io.BytesIO(data))
        chunks: List[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        chunks.append(t)

        text = "\n\n".join(chunks).strip()
        return text or None


register_extractor(PptxExtractor())
